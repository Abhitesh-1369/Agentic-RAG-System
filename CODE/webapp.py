import streamlit as st
import fitz
from docx import Document
from pptx import Presentation
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
import faiss
import cohere
import os
import time
import uuid

def parse_pdf(file_content):
    doc = fitz.open(stream=file_content, filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def parse_docx(file_content):
    doc = Document(file_content)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_content):
    prs = Presentation(file_content)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_csv(file_content):
    df = pd.read_csv(file_content)
    return df.to_string(index=False)

def parse_txt(file_content):
    return file_content.read().decode("utf-8")

def chunk_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def get_embeddings(text_chunks):
    vectorizer = TfidfVectorizer(max_features=384)
    vectors = vectorizer.fit_transform(text_chunks).toarray()
    if vectors.shape[1] < 384:
        padded = np.zeros((vectors.shape[0], 384))
        padded[:, :vectors.shape[1]] = vectors
        return padded
    return vectors

class VectorStore:
    def __init__(self, dim):
        self.index = faiss.IndexFlatL2(dim)
        self.chunks = []

    def add_embeddings(self, embeddings, chunks):
        self.index.add(np.array(embeddings))
        self.chunks.extend(chunks)

    def search(self, query_embedding, top_k=3):
        D, I = self.index.search(np.array([query_embedding]), top_k)
        return [self.chunks[i] for i in I[0]]

class LLMResponseAgent:
    def __init__(self, cohere_client):
        self.client = cohere_client

    def generate_response(self, context, query):
        prompt = f"Based on the following context, please answer the question.\n\nContext:\n{context}\n\nQuestion: {query}\n\nAnswer:"
        try:
            response = self.client.generate(
                model="command-r-plus",
                prompt=prompt,
                max_tokens=500,
                temperature=0.2
            )
            if response.generations:
                return response.generations[0].text.strip()
        except Exception as e:
            st.error(f"Error communicating with the LLM: {e}")
        return "Sorry, I couldn't generate a response."


st.set_page_config(
    page_title="Agentic RAG System",
    page_icon="🧠",
    layout="centered"
)

st.markdown("""
<style>
    /* Set a white background and black text for the entire app */
    .stApp {
        background-color: #FFFFFF;
        color: #000000;
    }
    /* Ensure all text elements inherit the black color */
    body, p, div, input, label, h1, h2, h3, h4, h5, h6 {
       color: #000000;
    }
    /* Style for individual chat messages with a black outline */
    .stChatMessage {
        background-color: #F0F2F6; /* Light grey background for the box */
        border: 1px solid #000000;   /* Black outline */
        border-radius: 0.5rem;
        padding: 1rem;
        margin-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)


if "messages" not in st.session_state:
    st.session_state.messages = []
if "vector_store" not in st.session_state:
    st.session_state.vector_store = VectorStore(dim=384)
if "file_processed" not in st.session_state:
    st.session_state.file_processed = False
if "llm_agent" not in st.session_state:
    try:
        cohere_api_key = "yAMTEoWRr6RWBkgw9czohN7YAVABVLcyltfXayfV"
        co_client = cohere.Client(cohere_api_key)
        st.session_state.llm_agent = LLMResponseAgent(co_client)
    except Exception as e:
        st.error("Could not initialize Cohere client. Please check your API key.")
        st.stop()


with st.sidebar:
    st.header("1. Upload Your Document")
    st.write("Upload a file to begin. The system will process it, enabling you to ask questions in the chat.")

    uploaded_file = st.file_uploader(
        "Choose a file",
        type=["pdf", "pptx", "docx", "csv", "txt"],
        label_visibility="collapsed"
    )

    if uploaded_file is not None:
        file_type = uploaded_file.name.split(".")[-1]
        file_content = uploaded_file

        parsers = {
            "pdf": parse_pdf, "pptx": parse_pptx, "docx": parse_docx,
            "csv": parse_csv, "txt": parse_txt
        }

        with st.spinner(f"Processing '{uploaded_file.name}'..."):
            st.session_state.vector_store = VectorStore(dim=384)
            st.session_state.messages = []

            parser_func = parsers.get(file_type)
            if parser_func:
                text = parser_func(file_content)
                chunks = chunk_text(text)
                embeddings = get_embeddings(chunks)
                st.session_state.vector_store.add_embeddings(embeddings, chunks)
                st.session_state.file_processed = True
                st.success("File processed successfully!")
                st.info("You can now ask questions in the main chat window.")
            else:
                st.error(f"Unsupported file type: {file_type}")
                st.session_state.file_processed = False


st.title("Agentic RAG System")
st.write("Welcome! Please upload a document in the sidebar to get started.")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask a question about your document...", disabled=not st.session_state.file_processed):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            query_embedding = get_embeddings([prompt])[0]
            top_chunks = st.session_state.vector_store.search(query_embedding)
            response = st.session_state.llm_agent.generate_response(top_chunks, prompt)
            st.markdown(response)

    st.session_state.messages.append({"role": "assistant", "content": response})
